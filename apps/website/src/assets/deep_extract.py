import fitz
import pptx
import docx
import os
import json

base = r'c:\Users\xchen\OneDrive\UU HCI Master'
inv_path = os.path.join(base, 'course_inventory.json')

with open(inv_path, 'r', encoding='utf-8') as f:
    courses = json.load(f)

for course_name, data in courses.items():
    c_path = data['path']
    course_text = f"--- {course_name} ({data['year_period']}) ---\n\n"
    
    count = 0
    for file_rel in data['files']:
        f_path = os.path.join(c_path, file_rel)
        lower_rel = file_rel.lower()
        
        # Skip textbooks, published papers, examples
        if 'textbook' in lower_rel or 'papers\\' in lower_rel or 'literature' in lower_rel or 'example' in lower_rel:
            continue
            
        ext = os.path.splitext(f_path)[1].lower()
        if ext not in ['.pdf', '.pptx', '.docx', '.md', '.txt']:
            continue
            
        if count > 15: # Extract more detail from the top 15 files of the course
            break

        course_text += f"\n\n==================================\n"
        course_text += f"FILE: {file_rel}\n"
        course_text += f"==================================\n\n"
        
        try:
            if ext == '.pdf':
                doc = fitz.open(f_path)
                # extract first 20 pages
                num_pages = min(20, len(doc))
                for i in range(num_pages):
                    course_text += f"--- Page {i+1} ---\n"
                    course_text += doc[i].get_text() + "\n"
                count += 1
                
            elif ext == '.pptx':
                prs = pptx.Presentation(f_path)
                num_slides = min(20, len(prs.slides))
                for i, slide in enumerate(list(prs.slides)[:num_slides]):
                    course_text += f"--- Slide {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            course_text += shape.text + "\n"
                count += 1
                
            elif ext == '.docx':
                d = docx.Document(f_path)
                course_text += "--- Document Content ---\n"
                for para in d.paragraphs[:100]:
                    course_text += para.text + "\n"
                count += 1
                
            elif ext in ['.txt', '.md']:
                with open(f_path, 'r', encoding='utf-8', errors='ignore') as f:
                    course_text += f.read(5000)
                count += 1
                
        except Exception as e:
            course_text += f"Error extracting: {e}\n"
            
    safe_name = course_name.replace(' ', '_').replace('/', '_')
    out_file = os.path.join(base, f'deep_{safe_name}.txt')
    with open(out_file, 'w', encoding='utf-8') as f:
        f.write(course_text)
        
print("Deep extraction completed for all courses.")
